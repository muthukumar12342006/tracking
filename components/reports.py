"""Executive report generation – PDF, Excel, PowerPoint-style exports."""

import io
from datetime import datetime
from pathlib import Path

import pandas as pd
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_LEFT
from reportlab.lib.pagesizes import A4, landscape
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import Image, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

from components.styles import COLORS, LOGO_PATH


def _get_logo_path() -> str | None:
    if LOGO_PATH.exists():
        return str(LOGO_PATH)
    return None


def generate_executive_pdf(summary: dict, sections: list[dict]) -> bytes:
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=A4, topMargin=0.75 * inch, bottomMargin=0.75 * inch)
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle("CapTitle", parent=styles["Heading1"], fontSize=22, textColor=colors.HexColor(COLORS["navy"]), spaceAfter=12)
    heading_style = ParagraphStyle("CapHeading", parent=styles["Heading2"], fontSize=14, textColor=colors.HexColor(COLORS["teal_dark"]), spaceBefore=16, spaceAfter=8)
    body_style = ParagraphStyle("CapBody", parent=styles["Normal"], fontSize=10, leading=14, textColor=colors.HexColor("#334155"))
    meta_style = ParagraphStyle("CapMeta", parent=styles["Normal"], fontSize=9, textColor=colors.HexColor("#64748B"), alignment=TA_CENTER)

    story = []
    logo = _get_logo_path()
    if logo:
        story.append(Image(logo, width=1.8 * inch, height=0.6 * inch))
        story.append(Spacer(1, 0.2 * inch))

    story.append(Paragraph("Cap AI – Executive Budget Report", title_style))
    story.append(Paragraph(f"Report Date: {datetime.now().strftime('%B %d, %Y')}", meta_style))
    story.append(Spacer(1, 0.3 * inch))

    story.append(Paragraph("Executive Summary", heading_style))
    for key, val in summary.items():
        story.append(Paragraph(f"<b>{key}:</b> {val}", body_style))

    story.append(Spacer(1, 0.2 * inch))
    story.append(Paragraph("Key Findings & Recommendations", heading_style))

    for section in sections:
        story.append(Paragraph(section.get("title", "Section"), heading_style))
        for item in section.get("items", []):
            story.append(Paragraph(f"• {item}", body_style))
        if section.get("table") is not None and not section["table"].empty:
            tbl_data = [list(section["table"].columns)] + section["table"].head(10).values.tolist()
            t = Table(tbl_data, repeatRows=1)
            t.setStyle(TableStyle([
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor(COLORS["navy"])),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("FONTSIZE", (0, 0), (-1, -1), 8),
                ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor(COLORS["gray_200"])),
                ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#F8FAFC")]),
            ]))
            story.append(t)
            story.append(Spacer(1, 0.15 * inch))

    doc.build(story)
    buffer.seek(0)
    return buffer.getvalue()


def export_dataframe_excel(df: pd.DataFrame, sheet_name: str = "Report") -> bytes:
    buffer = io.BytesIO()
    with pd.ExcelWriter(buffer, engine="xlsxwriter") as writer:
        df.to_excel(writer, sheet_name=sheet_name, index=False)
        workbook = writer.book
        worksheet = writer.sheets[sheet_name]
        header_fmt = workbook.add_format({"bold": True, "bg_color": COLORS["navy"], "font_color": "white", "border": 1})
        for col_num, value in enumerate(df.columns):
            worksheet.write(0, col_num, value, header_fmt)
            worksheet.set_column(col_num, col_num, max(12, len(str(value)) + 4))
    buffer.seek(0)
    return buffer.getvalue()


def build_executive_summary(variance_df, approval_df, overspend_df, rebudget_df, assumption_df) -> dict:
    summary = {}
    if variance_df is not None and not variance_df.empty:
        total_budget = variance_df["Budget Amount"].sum()
        total_actual = variance_df["Actual Amount"].sum()
        variance = total_actual - total_budget
        high_risk = len(variance_df[variance_df.get("Risk Rating", pd.Series()) == "High"]) if "Risk Rating" in variance_df.columns else 0
        summary["Total Budget"] = f"${total_budget:,.0f}"
        summary["Total Actual Spend"] = f"${total_actual:,.0f}"
        summary["Overall Variance"] = f"${variance:,.0f} ({variance/total_budget*100:+.1f}%)" if total_budget else "N/A"
        summary["High-Risk Variances"] = str(high_risk)
    else:
        summary["Status"] = "Load variance data for full executive summary"

    if approval_df is not None and not approval_df.empty and "Compliance Status" in approval_df.columns:
        critical = len(approval_df[approval_df["Compliance Status"] == "Critical"])
        summary["Approval Exceptions"] = str(critical)

    if overspend_df is not None and not overspend_df.empty and "Risk Level" in overspend_df.columns:
        chronic = len(overspend_df[overspend_df["Risk Level"].isin(["High", "Critical"])])
        summary["Chronic Overspend Heads"] = str(chronic)

    return summary


def build_report_sections(variance_df, approval_df, overspend_df, rebudget_df, assumption_df) -> list[dict]:
    sections = []

    if variance_df is not None and not variance_df.empty and "Risk Rating" in variance_df.columns:
        high = variance_df[variance_df["Risk Rating"] == "High"]
        sections.append({
            "title": "Significant Variances",
            "items": [f"{r['Budget Head']}: {r.get('Variance %', 0):+.1f}% variance" for _, r in high.head(5).iterrows()] or ["No high-risk variances detected"],
            "table": high[["Budget Head", "Department", "Variance Amount", "Variance %", "Risk Rating"]].head(10) if not high.empty else pd.DataFrame(),
        })

    if approval_df is not None and not approval_df.empty:
        issues = approval_df[approval_df.get("Compliance Status", "") != "Compliant"] if "Compliance Status" in approval_df.columns else pd.DataFrame()
        sections.append({
            "title": "Approval Exceptions",
            "items": [f"{r['Budget Head']}: {r.get('Compliance Status', 'Unknown')}" for _, r in issues.head(5).iterrows()] or ["All approvals compliant"],
        })

    if overspend_df is not None and not overspend_df.empty:
        sections.append({
            "title": "Overspending Trends",
            "items": [f"{r['Budget Head']}: {r.get('Overspend Frequency', 0)} breaches, Risk={r.get('Risk Level', 'N/A')}" for _, r in overspend_df.head(5).iterrows()],
        })

    sections.append({
        "title": "Executive Recommendations",
        "items": [
            "Review all High and Critical risk budget heads within 5 business days.",
            "Escalate missing approval records to Finance Controller immediately.",
            "Implement monthly variance review cadence for departments exceeding 10% variance.",
            "Validate assumption models against updated historical benchmarks quarterly.",
            "Require dual approval for all budget revisions exceeding 5% of original allocation.",
        ],
    })

    return sections


def generate_pptx_dashboard(summary: dict, kpi_data: list[tuple]) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Cap AI – Executive Dashboard"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(10, 22, 40)

    logo = _get_logo_path()
    if logo:
        slide.shapes.add_picture(logo, Inches(10.5), Inches(0.2), height=Inches(0.7))

    y = 1.5
    for i, (label, value) in enumerate(kpi_data[:6]):
        col = i % 3
        row = i // 3
        box = slide.shapes.add_textbox(Inches(0.5 + col * 4.2), Inches(y + row * 2.2), Inches(3.8), Inches(1.8))
        tf = box.text_frame
        tf.paragraphs[0].text = value
        tf.paragraphs[0].font.size = Pt(24)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(0, 168, 136)
        p2 = tf.add_paragraph()
        p2.text = label
        p2.font.size = Pt(12)
        p2.font.color.rgb = RGBColor(100, 116, 139)

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()
